import os
import tempfile
import streamlit as st
from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_chroma import Chroma
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser
from langchain_core.documents import Document
from datetime import datetime

# Word and PowerPoint loaders
from docx import Document as DocxDocument
from pptx import Presentation

load_dotenv()

st.set_page_config(page_title="📘 SHU Assistant", layout="wide", initial_sidebar_state="expanded")

# --- Custom CSS ---
st.markdown("""
<style>
    .stApp { background-color: #0f1117; }

    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1a1f2e 0%, #162032 100%);
        border-right: 1px solid #2d3748;
    }

    .main-title {
        background: linear-gradient(90deg, #4f8ef7, #a78bfa);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        font-size: 2.2rem;
        font-weight: 800;
        margin-bottom: 0;
    }

    .mode-badge {
        display: inline-block;
        padding: 4px 14px;
        border-radius: 20px;
        font-size: 0.78rem;
        font-weight: 600;
        margin-bottom: 10px;
    }
    .mode-pdf    { background: #1e3a5f; color: #63b3ed; border: 1px solid #2b6cb0; }
    .mode-word   { background: #1e3a28; color: #68d391; border: 1px solid #276749; }
    .mode-ppt    { background: #3a1e1e; color: #fc8181; border: 1px solid #9b2c2c; }
    .mode-general{ background: #1a2e1a; color: #68d391; border: 1px solid #276749; }

    .chip-container { display: flex; flex-wrap: wrap; gap: 8px; margin: 10px 0; }
    .chip {
        background: #1e2a3a;
        border: 1px solid #2d3f55;
        color: #90cdf4;
        padding: 6px 14px;
        border-radius: 20px;
        font-size: 0.82rem;
        cursor: pointer;
    }

    .stat-card {
        background: #1a2035;
        border: 1px solid #2d3748;
        border-radius: 10px;
        padding: 10px 14px;
        text-align: center;
        margin: 4px 0;
    }
    .stat-number { font-size: 1.4rem; font-weight: 700; color: #4f8ef7; }
    .stat-label  { font-size: 0.72rem; color: #718096; }

    .typing { color: #a0aec0; font-style: italic; font-size: 0.85rem; }

    .source-box {
        background: #1a2035;
        border-left: 3px solid #4f8ef7;
        padding: 8px 12px;
        border-radius: 0 8px 8px 0;
        font-size: 0.8rem;
        color: #90cdf4;
        margin-top: 8px;
    }

    .welcome-banner {
        background: linear-gradient(135deg, #1a1f3a 0%, #1a2a3a 100%);
        border: 1px solid #2d4a6a;
        border-radius: 12px;
        padding: 16px 20px;
        margin-bottom: 20px;
    }
</style>
""", unsafe_allow_html=True)

# --- Helpers to extract text from Word and PPT ---

def load_docx(file_path: str) -> list[Document]:
    """Extract paragraphs from a .docx file as LangChain Documents."""
    doc = DocxDocument(file_path)
    docs = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:
            docs.append(Document(page_content=text, metadata={"paragraph": i, "source": file_path}))
    return docs


def load_pptx(file_path: str) -> list[Document]:
    """Extract slide text from a .pptx file as LangChain Documents."""
    prs = Presentation(file_path)
    docs = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            docs.append(Document(
                page_content="\n".join(texts),
                metadata={"slide": slide_num, "source": file_path}
            ))
    return docs


def get_file_type_label(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower()
    return {"pdf": "PDF", "docx": "Word", "doc": "Word", "pptx": "PowerPoint", "ppt": "PowerPoint"}.get(ext, "Document")


# --- Session State Init ---
defaults = {
    "chat_history": [],
    "pdf_processed": False,
    "retriever": None,
    "doc_name": None,
    "doc_type": None,
    "total_questions": 0,
    "session_start": datetime.now().strftime("%H:%M"),
    "suggested_click": None,
    "tone": "Helpful & Friendly",
}
for k, v in defaults.items():
    if k not in st.session_state:
        st.session_state[k] = v

# --- Sidebar ---
with st.sidebar:
    st.markdown("<h2 style='color:#4f8ef7; margin-bottom:4px;'>📘 SHU Assistant</h2>", unsafe_allow_html=True)
    st.caption("Powered by Google Gemini")
    st.markdown("---")

    col1, col2 = st.columns(2)
    with col1:
        st.markdown(f"""
        <div class='stat-card'>
            <div class='stat-number'>{st.session_state['total_questions']}</div>
            <div class='stat-label'>Questions</div>
        </div>""", unsafe_allow_html=True)
    with col2:
        st.markdown(f"""
        <div class='stat-card'>
            <div class='stat-number'>{st.session_state['session_start']}</div>
            <div class='stat-label'>Session Start</div>
        </div>""", unsafe_allow_html=True)

    st.markdown("---")

    st.subheader("🎨 Response Tone")
    st.session_state["tone"] = st.selectbox(
        "Choose tone:",
        ["Helpful & Friendly", "Formal & Academic", "Simple & Brief", "Detailed & Thorough"],
        label_visibility="collapsed"
    )

    st.markdown("---")

    # ---- File Upload (PDF + Word + PPT) ----
    st.subheader("📂 Upload Document")
    st.caption("Supports PDF, Word (.docx), and PowerPoint (.pptx)")
    uploaded_file = st.file_uploader(
        "Upload a document",
        type=["pdf", "docx", "doc", "pptx", "ppt"],
        label_visibility="collapsed"
    )

    if uploaded_file:
        if st.button("📥 Process Document", use_container_width=True):
            st.session_state["pdf_processed"] = False
            st.session_state["retriever"] = None

            with st.spinner("Processing document..."):
                ext = uploaded_file.name.rsplit(".", 1)[-1].lower()

                # Save to a temp file
                suffix = f".{ext}"
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                    tmp.write(uploaded_file.getbuffer())
                    tmp_path = tmp.name

                # Load into LangChain Documents
                if ext == "pdf":
                    loader = PyPDFLoader(tmp_path)
                    documents = loader.load()
                elif ext in ("docx", "doc"):
                    documents = load_docx(tmp_path)
                elif ext in ("pptx", "ppt"):
                    documents = load_pptx(tmp_path)
                else:
                    st.error("Unsupported file type.")
                    st.stop()

                if not documents:
                    st.error("No readable text found in the document.")
                    st.stop()

                text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
                chunks = text_splitter.split_documents(documents)

                CHROMA_DB_DIR = "chroma_db_doc"
                os.makedirs(CHROMA_DB_DIR, exist_ok=True)

                embeddings_model = GoogleGenerativeAIEmbeddings(model="models/gemini-embedding-001")
                vector_store = Chroma.from_documents(
                    documents=chunks,
                    embedding=embeddings_model,
                    persist_directory=CHROMA_DB_DIR
                )

                st.session_state["retriever"]     = vector_store.as_retriever(search_kwargs={"k": 4})
                st.session_state["pdf_processed"] = True
                st.session_state["doc_name"]      = uploaded_file.name
                st.session_state["doc_type"]      = get_file_type_label(uploaded_file.name)
                st.session_state["pdf_chunks"]    = len(chunks)

                # Clean up temp file
                os.unlink(tmp_path)

            st.success(f"✅ Done! {len(chunks)} chunks indexed.")

    if st.session_state["pdf_processed"]:
        doc_type = st.session_state.get("doc_type", "Document")
        icon = {"PDF": "📄", "Word": "📝", "PowerPoint": "📊"}.get(doc_type, "📁")
        st.markdown(f"""
        <div class='source-box'>
            {icon} <b>{st.session_state['doc_name']}</b><br>
            🏷️ {doc_type} &nbsp;|&nbsp; 🔢 {st.session_state.get('pdf_chunks', '?')} chunks indexed
        </div>""", unsafe_allow_html=True)

    st.markdown("---")

    # Mode badge
    if st.session_state["pdf_processed"]:
        doc_type = st.session_state.get("doc_type", "Document")
        badge_class = {"PDF": "mode-pdf", "Word": "mode-word", "PowerPoint": "mode-ppt"}.get(doc_type, "mode-pdf")
        icon = {"PDF": "📄", "Word": "📝", "PowerPoint": "📊"}.get(doc_type, "📁")
        st.markdown(f"<div class='mode-badge {badge_class}'>{icon} {doc_type} Q&A Mode</div>", unsafe_allow_html=True)
    else:
        st.markdown("<div class='mode-badge mode-general'>💬 General Chat Mode</div>", unsafe_allow_html=True)

    st.markdown("---")

    if st.button("🗑️ Clear Chat", use_container_width=True):
        st.session_state["chat_history"] = []
        st.session_state["total_questions"] = 0
        st.rerun()

    if st.session_state["pdf_processed"]:
        if st.button("❌ Remove Document", use_container_width=True):
            st.session_state["pdf_processed"] = False
            st.session_state["retriever"]     = None
            st.session_state["doc_name"]      = None
            st.session_state["doc_type"]      = None
            st.rerun()

# --- Main Area ---
if not os.getenv("GOOGLE_API_KEY"):
    st.error("❌ Missing GOOGLE_API_KEY in your .env file.")
    st.stop()

st.markdown("<div class='main-title'>📘 Salim Habib University Assistant</div>", unsafe_allow_html=True)
st.caption("Ask anything, or upload a PDF / Word / PowerPoint file to chat with your documents!")

if not st.session_state["chat_history"]:
    st.markdown("""
    <div class='welcome-banner'>
        <b style='color:#4f8ef7;'>👋 Welcome!</b><br>
        <span style='color:#a0aec0; font-size:0.9rem;'>
        I'm your SHU AI Assistant. Chat freely, or upload a <b>PDF</b>, <b>Word (.docx)</b>,
        or <b>PowerPoint (.pptx)</b> file from the sidebar to ask questions about your documents.
        </span>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("**💡 Try asking:**")
    suggestions = [
        "What courses are offered at SHU?",
        "How do I apply for admission?",
        "What are the exam rules?",
        "Explain machine learning simply",
        "Help me write a study plan",
        "What is data science?",
    ]
    cols = st.columns(3)
    for i, suggestion in enumerate(suggestions):
        with cols[i % 3]:
            if st.button(suggestion, key=f"sug_{i}", use_container_width=True):
                st.session_state["suggested_click"] = suggestion
                st.rerun()

# --- Display chat history ---
for role, msg in st.session_state["chat_history"]:
    if role == "user":
        with st.chat_message("user", avatar="🧑‍🎓"):
            st.markdown(msg)
    elif role == "bot":
        with st.chat_message("assistant", avatar="🤖"):
            st.markdown(msg)
    elif role == "source":
        with st.expander("📚 View Sources Used", expanded=False):
            st.markdown(msg)

# --- Tone map ---
tone_map = {
    "Helpful & Friendly":   "Be warm, friendly, and encouraging in your response.",
    "Formal & Academic":    "Use formal, academic language suitable for university students.",
    "Simple & Brief":       "Give a very short and simple answer, easy to understand.",
    "Detailed & Thorough":  "Give a comprehensive, detailed, and well-structured answer.",
}

# --- Handle input ---
query = st.chat_input("💬 Ask me anything...")
if st.session_state["suggested_click"]:
    query = st.session_state["suggested_click"]
    st.session_state["suggested_click"] = None

if query:
    st.session_state["total_questions"] += 1

    with st.chat_message("user", avatar="🧑‍🎓"):
        st.markdown(query)
    st.session_state["chat_history"].append(("user", query))

    with st.chat_message("assistant", avatar="🤖"):
        with st.spinner("Thinking... 🤔"):
            llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.3)
            tone_instruction = tone_map[st.session_state["tone"]]
            doc_type = st.session_state.get("doc_type", "document")

            # --- Document Q&A Mode ---
            if st.session_state.get("pdf_processed") and st.session_state.get("retriever"):
                retriever = st.session_state["retriever"]

                prompt = ChatPromptTemplate.from_template(f"""
                You are an intelligent assistant for Salim Habib University students.
                {tone_instruction}
                Use the context below from the uploaded {doc_type} to answer the question.
                If the answer is not in the context, clearly say so.
                At the end, mention which part of the document helped you answer.

                Context: {{context}}
                Question: {{input}}
                """)

                def format_docs(docs):
                    formatted = []
                    for d in docs:
                        meta = d.metadata
                        # Label source location based on doc type
                        if "page" in meta:
                            loc = f"Page {meta['page'] + 1}"
                        elif "slide" in meta:
                            loc = f"Slide {meta['slide']}"
                        elif "paragraph" in meta:
                            loc = f"Paragraph {meta['paragraph'] + 1}"
                        else:
                            loc = "Section"
                        formatted.append(f"[{loc}]: {d.page_content}")
                    return "\n\n".join(formatted)

                chain = (
                    {"context": retriever | format_docs, "input": RunnablePassthrough()}
                    | prompt
                    | llm
                    | StrOutputParser()
                )

                answer = chain.invoke(query)

                # Build source preview
                source_docs = retriever.invoke(query)
                source_text = ""
                for i, doc in enumerate(source_docs[:3]):
                    meta = doc.metadata
                    if "page" in meta:
                        loc = f"Page {int(meta['page']) + 1}"
                    elif "slide" in meta:
                        loc = f"Slide {meta['slide']}"
                    elif "paragraph" in meta:
                        loc = f"Paragraph {meta['paragraph'] + 1}"
                    else:
                        loc = "Section"
                    preview = doc.page_content[:200].strip()
                    source_text += f"**Source {i+1} — {loc}:**\n> {preview}...\n\n"

                st.session_state["chat_history"].append(("source", source_text))

            # --- General Chat Mode ---
            else:
                prompt = ChatPromptTemplate.from_template(f"""
                You are a helpful AI assistant for Salim Habib University students.
                {tone_instruction}
                Answer the following question clearly.

                Question: {{question}}
                """)

                chain = prompt | llm | StrOutputParser()
                answer = chain.invoke({"question": query})

        st.markdown(answer)
        st.caption(f"🕐 {datetime.now().strftime('%H:%M:%S')}")

    st.session_state["chat_history"].append(("bot", answer))
    st.rerun()